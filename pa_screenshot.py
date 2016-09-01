# -*- coding:utf-8 -*-
import sys, requests, time
import pandas as pd
import webbrowser as wb
import pyscreenshot
from PIL import Image, ImageOps
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches, Pt
import selenium
from selenium import  webdriver
import validators

# driver = webdriver.Chrome('/Users/apple/Downloads/chromedriver')
# driver.set_window_size(2000, 1440)

# 自动截图并裁剪

xls = pd.ExcelFile('links.xlsx')
df = pd.read_excel(xls, header = None, index_col = None, na_value = None)

rows_df = len(df.index.values.tolist())
bad_list = []
# 创建 Word 和 PPT 文档
document = Document()
document.add_heading('Document Title', 0)

prs = Presentation()
blank_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(blank_slide_layout)



for i in range(1, rows_df):
    wb.open(df.ix[i, 5],new = 0)
    time.sleep(10)
    r = requests.head(df.ix[i, 5])
    if r.status_code == 200:
        img = pyscreenshot.grab()
    else:
        bad_list.append(i)
        continue
    img2 = img.crop((0,240,2000,1440))
    img3 = ImageOps.expand(img2, border = 10, fill = 'black') # 加黑框
    img3.save('{}.png'.format(i))

    # driver.close()

    # 写入 Word 文档
    headline = df.ix[i, 3]
    p = document.add_heading('{}'.format(headline), level = 1)
    document.add_picture('{}.png'.format(i))
    document.save('report.docx')

    # 生成 PPT
    shapes = slide.shapes

    # 增加 PPT 每页标题
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = '海外重量级网站传播情况: {}'.format(df.ix[i, 3])
    p.font.bold = True
    p.font.size = Pt(30)


    # 插入截图
    top = Inches(2.5)
    left = Inches(1.5)
    height = Inches(4)
    pic = slide.shapes.add_picture('{}.png'.format(i), left, top, height=height)
    prs.save('report.pptx')
    slide = prs.slides.add_slide(blank_slide_layout)




if bad_list != []:
    print('''These links are inaccessible: ''', bad_list, ''' Please check them again.
     ----------THE END----------''')
else:
    pass

print('''----------------------------------------
I have fought the good fight,           |
I have finished the race,               |
I have kept the faith.                  |
                        2 Timothy 4:7   |
----------------------------------------


Have a good day.
: ) ''')