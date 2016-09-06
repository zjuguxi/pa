# -*- coding:utf-8 -*-
import os
import pandas as pd
import xlsxwriter
import sys, requests, time
import webbrowser as wb
import pyscreenshot
from PIL import Image, ImageOps
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches, Pt
import selenium
from selenium import  webdriver
import validators

# 得到当前path
cwd = os.getcwd()

progress_bar = 0 # 全局变量 progress_bar

# 打开 Excel 文档的 Pickup 汇总表，读取链接
xls = pd.ExcelFile('1.xls')
df_pickup = pd.read_excel(xls, 'Pickup', header = None, index_col = None, na_value = None)

# Pickup 的总 row 数
p = len(df_pickup.index.values.tolist())

# 找到 Pickup column[0] 里 Story Number 第一次出现的位置，取其 index，命名为 m
# 从 Story Number 处开始截取 DataFrame
s = df_pickup.iloc[:, 0]
m = s[s == 'Story Number'].index[0]
df = pd.read_excel(xls, 'Pickup', header = m, index_col = None, na_value = None)

# 找到 Twitter 的 index，命名为 twt
# 从 Twitter Handle 处开始截取 DataFrame
twt = s[s == 'Twitter Handle'].index[0]
df_twitter = pd.read_excel(xls, 'Pickup', header = twt, index_col = None, na_value = None)
# 把 Twitter 信息生成新表格
writer = pd.ExcelWriter('twitter.xlsx', engine='xlsxwriter')
df_twitter.to_excel(writer, sheet_name='Sheet1')
writer.save()

# 打开 Realeases 读取文章标题
df_headline = pd.read_excel(xls, 'Releases', header = None, index_col = 0, na_value = None)

# Pickup 的总 row 数
p = len(df.index.values.tolist())
# Releases 的总 row 数
n = len(df_headline.index.values.tolist())

# 提取 Releases 里对应的 Story Number和 Headline，形成 Series
series_headline = []
series_storynumber = []

for i in range(n):
    if df_headline.index[i] == 'Headline':
        series_headline.append(df_headline.ix[i, 1])
    if df_headline.index[i] == 'Story Number':
        series_storynumber.append(df_headline.ix[i, 1])

df_releases = pd.DataFrame(series_headline, index = series_storynumber) # 将2个series合并为一个dataframe

# Releases新表格 的总 row 数
q = len(df_releases.index.values.tolist())

print(df_releases)

list_addon = []

for i in range(p):
    for j in range(q):
        if df.ix[i, 'Story Number'] == df_releases.index[j]: # 对比新 Releases 和 Pickup 的Story Number
            df.ix[i, 'Headline'] = df_releases.ix[j, 0]

df = df[df.Headline.notnull()] # 去掉 Headline 为空的 rows，仅剩普通Links

writer = pd.ExcelWriter('links.xlsx', engine='xlsxwriter', options={'strings_to_urls': False})
df.to_excel(writer, sheet_name='Sheet1')
writer.save()

##### 以上为数据处理 #####


##### 以上为写入 DOC 和 PPT #####

xls = pd.ExcelFile('links.xlsx')
df = pd.read_excel(xls, header = None, index_col = None, na_value = None)

xls_twitter = pd.ExcelFile('twitter.xlsx')
df_twitter = pd.read_excel(xls_twitter, header = None, index_col = None, na_value = None)

rows_df = len(df.index.values.tolist())
rows_df_twitter = len(df_twitter.index.values.tolist())
bad_list = []
# 创建 Word 和 PPT 文档
document = Document()
document.add_heading('Document Title', 0)

prs = Presentation()
blank_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(blank_slide_layout)


for i in range(1, rows_df):
    progress_bar global
    progress_bar += 1
    print('''=========
        Working on this Now...............
            {}
            ========='''.format(df[df.index == i]))
    wb.open(df.ix[i, 5],new = 0)
    time.sleep(10)
    r = requests.head(df.ix[i, 5])
    if r.status_code == 200:
        img = pyscreenshot.grab()
    else:
        bad_list.append(i)
        continue
    img2 = img.crop((0,240,2000,1440))
    img3 = ImageOps.expand(img2, border = 10, fill = 'black') # 加黑框
    img3.save('{}.png'.format(i))

    # 写入 Word 文档
    headline = df.ix[i, 3]
    p = document.add_heading('{}'.format(headline), level = 1)
    document.add_picture('{}.png'.format(i))
    document.save('report.docx')

    # 生成 PPT
    shapes = slide.shapes

    # 增加 PPT 每页标题
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = '海外重量级网站传播情况: {}'.format(df.ix[i, 3])
    p.font.bold = True
    p.font.size = Pt(30)


    # 插入截图
    top = Inches(2.5)
    left = Inches(1.5)
    height = Inches(4)
    pic = slide.shapes.add_picture('{}.png'.format(i), left, top, height=height)
    prs.save('report.pptx')
    slide = prs.slides.add_slide(blank_slide_layout)

for i in range(1, rows_df_twitter):
    progress_bar global
    progress_bar += 1
    print('''=========
        Working on this Now......
            {}
            ========='''.format(df_twitter[df_twitter.index == i]))
    wb.open(df_twitter.ix[i, 2],new = 0)
    time.sleep(10)
    r = requests.head(df_twitter.ix[i, 2])
    if r.status_code == 200:
        img = pyscreenshot.grab()
    else:
        bad_list.append(i)
        continue
    img2 = img.crop((0,240,2000,1440))
    img3 = ImageOps.expand(img2, border = 10, fill = 'black') # 加黑框
    img3.save('{}.png'.format(i))

    # 写入 Word 文档
    headline = 'Twitter 账号: ', df_twitter.ix[i, 1]
    p = document.add_heading('{}'.format(headline), level = 1)
    document.add_picture('{}.png'.format(i))
    document.save('report.docx')

    # 生成 PPT
    shapes = slide.shapes

    # 增加 PPT 每页标题
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = 'Twitter 账号: {}'.format(df_twitter.ix[i, 1])
    p.font.bold = True
    p.font.size = Pt(30)


    # 插入截图
    top = Inches(2.5)
    left = Inches(1.5)
    height = Inches(4)
    pic = slide.shapes.add_picture('{}.png'.format(i), left, top, height=height)
    prs.save('report.pptx')
    slide = prs.slides.add_slide(blank_slide_layout)

if bad_list != []:
    print('''These links are inaccessible: ''', bad_list, ''' Please check them again.
     ----------THE END----------''')
else:
    pass

print('''----------------------------------------
I have fought the good fight,           |
I have finished the race,               |
I have kept the faith.                  |
                        2 Timothy 4:7   |
----------------------------------------


Have a good day.
: ) ''')